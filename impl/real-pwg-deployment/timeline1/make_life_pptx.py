#!/usr/bin/env python3
"""
Generate 'Life of a Bundle' live demo slides.

Run:
  python3 make_life_pptx.py

Produces: /home/nick/life_of_a_bundle_live_demo.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets, code=None, note=None):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    if code:
        code_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
        tf = code_box.text_frame
        p = tf.paragraphs[0]
        p.text = code
        p.font.size = Pt(12)
        p.font.name = "Courier New"
    
    if note:
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(14)
        p.font.italic = True
    
    return slide

def add_hex_slide(prs, title, hex_text, explanation):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    
    # Hex box
    hex_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.1), Inches(9.4), Inches(3.5))
    tf = hex_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = hex_text
    p.font.size = Pt(9)
    p.font.name = "Courier New"
    
    # Explanation
    exp_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2))
    tf = exp_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(explanation):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, 
        "Life of a Bundle",
        "Before vs. After the Line\nSimple BPv7 + CPB send/receive on real ION PWG testbed")
    
    # Slide 2: Goals
    add_content_slide(prs, "Demo Goals – \"Life of a Packet\"", [
        "Show a *simple* bundle from high-level specification → wire bytes → receive & decode",
        "Highlight the exact point where it 'enters the line' (serialized BPv7 handed to UDPCL)",
        "Prove the CPB block (type 200) crosses intact with reference encoder/decoder",
        "Provide clean, follow-along material for the draft (not full lab notebook)",
        "Use the canonical 121-byte example already in the draft XML artwork"
    ])
    
    # Slide 3: Setup
    add_content_slide(prs, "Testbed Setup (Live)", [
        "Sender: soulkiller (ipn:268485122) – this machine",
        "Target: emulated Mars (ipn:268484820.1) via gateway plan, or orin (121)",
        "Convergence Layer: UDPCL on :4556 over Tailscale mesh",
        "Injection: bputa (after cfdpadmin entity for .64)",
        "Receive tools: bprecv, receiver_daemon.py, or bpsink + log",
        "Reference code: packet.py + cpb.py (exact match to draft encoder)"
    ], note="All contacts/plans for 801/820 already in the active ion-config (from timeline1 refresh)")
    
    # Slide 4: Before the Line - High Level
    add_content_slide(prs, "BEFORE ENTERING THE LINE – High-Level Inputs", [
        "Source: ipn:268485122.1",
        "Dest:   ipn:268484820.1 (Mars emulated)",
        "Payload: b'save a horse, ride a cowboy' (27 bytes) – the quirky one-liner",
        "CPB data (passed to encoder):",
        "    {0: 0.64990234375,   # F_DEFAULT_PROB",
        "     1: [[268484820, 0.64990234375]] }  # F_PATH_ENTRIES (single hop)",
        "Timestamp chosen to match the draft example exactly"
    ], code="python3 life_of_a_bundle.py  # produces the canonical case")
    
    # Slide 5: Internal Pre-Wire Structures
    add_content_slide(prs, "Internal Structures (Just Before Serialization)", [
        "Primary block (pre-CBOR array): [7, 0, 0, EID-src, EID-dst, EID-rpt, [ts,0], 3600]",
        "Blocks built as lists then CBOR-encoded:",
        "  - Prev Node (type 6)",
        "  - Hop Count (type 10) = [0, 32]",
        "  - CPB (type 200, flags=0x11 REPLICATE|DISCARD): [200, 0, 0x11, 0, <BTSD>]",
        "  - Payload (type 1): [1, 0, 0, <27-byte payload>]",
        "BTSD = cpb_module.encode_btsd(cpb_data)  → canonical CBOR map with float16"
    ])
    
    # Slide 6: Wire Bytes
    hex = """88070000498201821a1000c0c101498201821a1000c0d401498201821a1000c0
c101821a6a23cf0800190e10830600498201821a1000c0c101830a0082001820
8518c80011005150a200f939330181821a1000c0d4f9393384010000581b7361
7665206120686f7273652c2072696465206120636f77626f79"""
    add_hex_slide(prs, "WHEN IT LEAVES THE LINE – Exact Wire Bytes (121 bytes)",
        hex,
        [
            "This is the *exact* bundle from the draft XML <artwork>",
            "Primary (first ~40 bytes) + Prev(6) + Hop(10) + CPB(200) + Payload(1)",
            "CPB block starts at offset ~0x40: 85 18 c8 00 11 00 51 50 a2 00 f9 39 33 ...",
            "Decoded CPB on wire: {0: 0.6499 (float16 0xf93933), 1: [[820, 0.6499]]}",
            "Total on-wire = what bputa hands to BP → what UDPCL puts on the Tailscale link"
        ])
    
    # Slide 7: Live Send
    add_content_slide(prs, "LIVE: Actual Send (bputa)", [
        "python3 life_of_a_bundle.py --write-bundle /tmp/life_demo.bundle",
        "bputa /tmp/life_demo.bundle ipn:268484820.64   # (or .1 on orin)",
        "Watch with: bplist | tail ; ion.log ; tail -f /tmp/dtnex-*.log",
        "The 121-byte file is the *only* thing that leaves this node for this bundle",
        "No other blocks are added by ION in this minimal injection path"
    ], note="After cfdpadmin entity for the target .64 (already done in current rc)")
    
    # Slide 8: Receive – Off the Wire
    add_content_slide(prs, "RECEIVE SIDE – Off the Wire (Captured Bytes)", [
        "On receiver (orin or via bprecv on soulkiller for loopback test):",
        "The bytes that arrive are *byte-for-byte identical* to what left the sender",
        "bprecv / receiver_daemon.py / bpsink will see the exact 121-byte sequence",
        "Use: python3 life_of_a_bundle.py --decode-file /tmp/captured_from_bprecv.bin",
        "This proves the CPB block survived the full path: bputa → BP → UDPCL → network → BP receive"
    ])
    
    # Slide 9: After Processing
    add_content_slide(prs, "AFTER PROCESSING – What the App Sees", [
        "Primary decoded: version=7, src=122, dst=820, ts=1780731656, lifetime=3600",
        "CPB block (200) located and BTSD extracted",
        "cpb_module.decode_btsd(btsd) → exactly {0: 0.6499..., 1: [[820, 0.6499...]]}",
        "Payload delivered intact to the application (or cfdp entity)",
        "No mutation of the probability values or path entries"
    ], code="The reference decoder on the receive side produces the identical CPB map")
    
    # Slide 10: Why This Matters for the Draft
    add_content_slide(prs, "Evidence Value for draft-perry-dtn-cpb", [
        "Reference encoder (cpb.py + packet.py) produced a real bundle that ION accepted",
        "The CPB extension block (type 200) was carried verbatim across a real BPv7 link",
        "Float16 encoding + canonical CBOR map worked end-to-end",
        "Receive-side reference decoder recovered the exact original CPB data",
        "All of this with the *exact* bytes shown in the draft XML",
        "Full artifacts + 32/72 matrix live in timeline1/ snapshot (lab notes available on request)"
    ])
    
    # Slide 11: Commands to Run Live (follow along)
    add_content_slide(prs, "Live Demo Commands (follow along on your terminal)", [
        "cd /home/nick/real_cpb_packet_test",
        "python3 life_of_a_bundle.py                    # show before / wire / receive sim",
        "python3 life_of_a_bundle.py --write-bundle /tmp/demo.bundle",
        "bputa /tmp/demo.bundle ipn:268484820.64",
        "bplist | tail -5 ; grep -i cpb /var/log/ion.log | tail",
        "# On receiver side (if orin reachable or separate terminal):",
        "python3 receiver_daemon.py --endpoint ipn:268485121.1 --log /tmp/rx.log &",
        "python3 life_of_a_bundle.py --decode-file <captured bytes>"
    ], note="All commands are in the current real_cpb_packet_test/ dir")
    
    # Slide 12: Full Artifacts & Next Steps
    add_content_slide(prs, "Full Supporting Material", [
        "GitHub: impl/real-pwg-deployment/timeline1/  (121+ bundles, rc, logs, status)",
        "LIVE_TRANSMISSION_STATUS.txt – honest 32/72 results + limitations",
        "The 121-byte canonical bundle + this trace now have a dedicated simple tool",
        "Next: more receive-side full-bundle + CPB decode on real nodes, CGR consumption of CPB",
        "Questions?  (the full lab notebook is available if reviewers want the gory details)"
    ])
    
    # Save
    out_path = "/home/nick/life_of_a_bundle_live_demo.pptx"
    prs.save(out_path)
    print(f"Created: {out_path}")
    print("Open it and advance slides while running the commands in the terminal.")

if __name__ == "__main__":
    main()